"""
PPT Component Engine — Structure Parser
========================================
Deconstructs .pptx templates slide-by-slide into PageEntry records
with 3-level tag classification.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .schema import (
    PageEntry, PageType, LayoutVariant, ElementMap,
    DesignTokens, TextSlot, ContentConstraints,
)
from .page_store import PageStore


class StructureParser:
    """
    Analyzes a .pptx template and produces PageEntry records for each slide.

    Classification rules (L2):
        cover          -> slide 0 + text_boxes <= 4
        ending         -> last slide + text_boxes <= 3
        toc            -> within first 3 slides + contains TOC keywords
        section-header -> not first/last + text_boxes 1-3 + has title
        image-full     -> 1 image + 0-1 text boxes
        image-left     -> image in left half
        image-right    -> image in right half
        chart          -> has chart placeholder
        table          -> has table
        content-3col   -> text_boxes >= 7
        content-2col   -> text_boxes >= 5
        content-1col   -> text_boxes 1-4 (default content)
    """

    TOC_KEYWORDS = ["目录", "contents", "agenda", "outline", "index", "目錄"]
    ENDING_KEYWORDS = ["感谢", "thanks", "thank you", "q&a", "谢", "致谢", "谢谢"]

    def parse_template(
        self, pptx_path: str | Path, family_name: str,
        page_store: Optional[PageStore] = None, commit: bool = True,
    ) -> list[PageEntry]:
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"Template not found: {pptx_path}")
        if pptx_path.suffix.lower() not in (".pptx",):
            raise ValueError(f"Expected .pptx, got: {pptx_path.suffix}")

        prs = Presentation(str(pptx_path))
        entries: list[PageEntry] = []
        total_slides = len(prs.slides)

        for idx, slide in enumerate(prs.slides):
            try:
                entry = self._parse_slide(slide, idx, total_slides, pptx_path.name, family_name)
                entries.append(entry)
                if page_store is not None:
                    page_store.add(entry)
            except Exception as e:
                warning_entry = self._error_entry(idx, pptx_path.name, family_name, str(e))
                entries.append(warning_entry)
                if page_store is not None:
                    page_store.add(warning_entry)

        if page_store is not None and commit:
            page_store.commit()
        return entries

    def parse_templates(
        self, templates: list[tuple[str | Path, str]], page_store: PageStore,
    ) -> dict[str, list[PageEntry]]:
        results = {}
        for pptx_path, family_name in templates:
            entries = self.parse_template(pptx_path, family_name, page_store, commit=False)
            results[family_name] = entries
        page_store.commit()
        return results

    def _parse_slide(self, slide, idx, total, source_name, family_name) -> PageEntry:
        element_map = self._build_element_map(slide)
        text_slots = self._extract_text_slots(slide)
        page_type = self._classify_page_type(idx, total, slide, element_map, text_slots)
        layout_variant = self._build_layout_variant(slide, element_map, text_slots)
        design_tokens = self._extract_design_tokens(slide, text_slots)
        content_constraints = self._compute_constraints(element_map, text_slots)
        quality = self._compute_quality(element_map, design_tokens, text_slots)
        warnings = []
        if design_tokens.incomplete:
            warnings.append("Some design tokens could not be extracted")
        if page_type == PageType.UNKNOWN:
            warnings.append("Page type could not be determined")

        return PageEntry(
            page_id=f"{family_name}_page_{idx:02d}",
            source_template=source_name,
            template_family=family_name,
            slide_index=idx,
            page_type=page_type,
            layout_variant=layout_variant,
            element_map=element_map,
            design_tokens=design_tokens,
            text_structure=text_slots,
            content_constraints=content_constraints,
            quality_score=quality,
            source_file="",
            warnings=warnings,
        )

    def _error_entry(self, idx, source_name, family_name, error) -> PageEntry:
        return PageEntry(
            page_id=f"{family_name}_page_{idx:02d}",
            source_template=source_name,
            template_family=family_name,
            slide_index=idx,
            page_type=PageType.UNKNOWN,
            quality_score=0.0,
            warnings=[f"Parse error: {error}"],
        )

    def _build_element_map(self, slide) -> ElementMap:
        em = ElementMap()
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                em.text_boxes += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                em.image_placeholders += 1
            if shape.has_table:
                em.tables += 1
            if shape.has_chart:
                em.chart_placeholders += 1
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                em.auto_shapes += 1
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                em.groups += 1
            if shape.shape_type == MSO_SHAPE_TYPE.LINE_CONNECTOR:
                em.connectors += 1
        return em

    def _extract_text_slots(self, slide) -> list[TextSlot]:
        slots = []
        slide_w = slide.slide_width or Emu(9144000)
        slide_h = slide.slide_height or Emu(6858000)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if not text:
                continue

            left = shape.left or 0
            top = shape.top or 0
            w = shape.width or 0
            h = shape.height or 0
            pos = self._compute_position(left, top, w, h, slide_w, slide_h)
            role = self._classify_text_role(shape, text, pos)

            font_name = ""
            font_size_pt = 0.0
            font_bold = False
            font_color = ""
            try:
                tf = shape.text_frame
                if tf.paragraphs:
                    for para in tf.paragraphs:
                        if para.runs:
                            run = para.runs[0]
                            font_name = run.font.name or ""
                            try:
                                font_size_pt = run.font.size.pt if run.font.size else 0.0
                            except Exception:
                                pass
                            font_bold = run.font.bold or False
                            try:
                                rgb = run.font.color.rgb
                                font_color = str(rgb) if rgb else ""
                            except Exception:
                                pass
                            break
            except Exception:
                pass

            if w and font_size_pt:
                char_width_emu = Pt(font_size_pt) * 0.6
                chars_per_line = max(1, int(w / char_width_emu))
            else:
                chars_per_line = 40

            if h and font_size_pt:
                line_height_emu = Pt(font_size_pt) * 1.3
                max_lines = max(1, int(h / line_height_emu))
            else:
                max_lines = 10

            max_chars = chars_per_line * max_lines

            alignment = ""
            try:
                if shape.text_frame.paragraphs:
                    al = shape.text_frame.paragraphs[0].alignment
                    alignment = str(al) if al else ""
            except Exception:
                pass

            slots.append(TextSlot(
                role=role, position=pos, font_size_pt=font_size_pt,
                max_chars=max_chars, font_name=font_name, font_bold=font_bold,
                font_color=font_color, alignment=alignment,
                shape_name=shape.name, shape_id=shape.shape_id,
            ))
        return slots

    def _compute_position(self, left, top, w, h, slide_w, slide_h) -> str:
        cx = (left + w / 2)
        cy = (top + h / 2)
        h_third = slide_w / 3
        v_third = slide_h / 3

        x_pos = "center"
        if cx < h_third:
            x_pos = "left"
        elif cx > h_third * 2:
            x_pos = "right"

        y_pos = "center"
        if cy < v_third:
            y_pos = "top"
        elif cy > v_third * 2:
            y_pos = "bottom"

        if x_pos == "center" and y_pos == "center":
            return "center"
        return f"{y_pos}-{x_pos}"

    def _classify_text_role(self, shape, text, pos) -> str:
        try:
            font_size = 0.0
            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                run = shape.text_frame.paragraphs[0].runs[0]
                font_size = run.font.size.pt if run.font.size else 0.0
        except Exception:
            font_size = 0.0

        if "top" in pos and font_size >= 24 and len(text) <= 60:
            return "title"
        if font_size >= 28:
            return "title"
        if "top" in pos and 16 <= font_size < 24:
            return "subtitle"
        if font_size <= 12 and len(text) <= 30:
            return "caption"
        if font_size <= 10:
            return "label"
        return "body"

    def _classify_page_type(self, idx, total, slide, em, slots) -> PageType:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text.lower())
        all_text_lower = " ".join(texts)

        if idx == 0 and em.text_boxes <= 4:
            return PageType.COVER

        if idx == total - 1 and em.text_boxes <= 3:
            if any(kw in all_text_lower for kw in self.ENDING_KEYWORDS):
                return PageType.ENDING
            if em.text_boxes <= 2:
                return PageType.ENDING

        if idx <= 2 and any(kw in all_text_lower for kw in self.TOC_KEYWORDS):
            return PageType.TOC

        if em.image_placeholders >= 1:
            if em.text_boxes <= 1:
                return PageType.IMAGE_FULL
            img_pos = self._get_first_image_position(slide)
            if img_pos == "left":
                return PageType.IMAGE_LEFT
            elif img_pos == "right":
                return PageType.IMAGE_RIGHT
            else:
                return PageType.IMAGE_FULL

        if em.chart_placeholders > 0:
            return PageType.CHART
        if em.tables > 0:
            return PageType.TABLE

        if em.text_boxes >= 7:
            return PageType.CONTENT_3COL
        if em.text_boxes >= 5:
            return PageType.CONTENT_2COL

        if 0 < idx < total - 1 and 1 <= em.text_boxes <= 3:
            has_title = any(s.role == "title" for s in slots)
            if has_title:
                return PageType.SECTION_HEADER

        if em.text_boxes == 1 and any("center" in s.position for s in slots):
            if slots and slots[0].font_size_pt >= 18:
                return PageType.QUOTE

        if em.text_boxes <= 4:
            return PageType.CONTENT_1COL

        return PageType.UNKNOWN

    def _get_first_image_position(self, slide) -> str:
        slide_w = slide.slide_width or Emu(9144000)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                cx = (shape.left or 0) + (shape.width or 0) / 2
                if cx < slide_w / 2:
                    return "left"
                return "right"
        return "unknown"

    def _build_layout_variant(self, slide, em, slots) -> LayoutVariant:
        lv = LayoutVariant()
        positions = set(s.position for s in slots if s.position)
        if len(positions) >= 3:
            lv.column_count = 3
        elif len(positions) >= 2:
            lv.column_count = 2
        else:
            lv.column_count = 1

        if em.image_placeholders > 0:
            img_pos = self._get_first_image_position(slide)
            lv.image_position = img_pos
            if img_pos == "left":
                lv.arrangement = "image-left-text-right"
            elif img_pos == "right":
                lv.arrangement = "text-left-image-right"
            else:
                lv.arrangement = "image-full"
        elif lv.column_count >= 2:
            lv.arrangement = "multi-column"
        else:
            lv.arrangement = "single-column"

        lv.has_title_bar = any(s.role == "title" for s in slots)
        lv.has_subtitle = any(s.role == "subtitle" for s in slots)
        text_positions = [s.position for s in slots if s.role in ("title", "body")]
        if text_positions:
            lv.text_position = max(set(text_positions), key=text_positions.count)
        return lv

    def _extract_design_tokens(self, slide, slots) -> DesignTokens:
        dt = DesignTokens()
        title_slots = [s for s in slots if s.role == "title"]
        body_slots = [s for s in slots if s.role == "body"]
        if title_slots:
            dt.font_title = title_slots[0].font_name
            dt.title_size_pt = title_slots[0].font_size_pt
        if body_slots:
            dt.font_body = body_slots[0].font_name
            dt.body_size_pt = body_slots[0].font_size_pt

        try:
            layout = slide.slide_layout
            bg = layout.background
            fill = bg.fill
            if fill.type is not None:
                try:
                    xml = fill._fill.xml
                    match = re.search(r'srgbClr val="([^"]+)"', xml)
                    if match:
                        dt.primary_color = f"#{match.group(1)}"
                except Exception:
                    pass
        except Exception:
            pass

        if not dt.primary_color:
            colors = [s.font_color for s in slots if s.font_color]
            if colors:
                dt.primary_color = max(set(colors), key=colors.count)

        dt.incomplete = not dt.primary_color or not dt.font_title
        return dt

    def _compute_constraints(self, em, slots) -> ContentConstraints:
        cc = ContentConstraints()
        for slot in slots:
            if slot.role == "title":
                cc.max_title_chars = max(cc.max_title_chars, slot.max_chars)
            elif slot.role == "subtitle":
                cc.max_subtitle_chars = max(cc.max_subtitle_chars, slot.max_chars)
            elif slot.role == "body":
                cc.max_body_lines = max(cc.max_body_lines, int(slot.max_chars / 40))
                cc.max_body_chars_per_line = max(cc.max_body_chars_per_line, 40)
        cc.requires_image = em.image_placeholders > 0
        cc.requires_chart = em.chart_placeholders > 0
        return cc

    def _compute_quality(self, em, dt, slots) -> float:
        score = 1.0
        if em.text_boxes == 0 and em.image_placeholders == 0:
            score = 0.0
        if dt.incomplete:
            score -= 0.2
        if not slots:
            score -= 0.1
        if dt.primary_color == "":
            score -= 0.1
        return max(0.0, score)
