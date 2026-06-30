"""
PPT Component Engine — Template Filler
========================================
4-phase pipeline: Match -> Analyze -> Assemble -> Validate.
Non-destructive: source .pptx files are never modified.
"""
from __future__ import annotations

import re
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .schema import (
    PageType, HarmonyMode, Severity,
    PageContent, ContentOutline, PageEntry,
    MatchResult, ContentFitResult, ContentIssue,
    AssemblyPlan, FillResult, ValidationIssue,
)
from .page_store import PageStore


class PageMatcher:
    def __init__(self, page_store: PageStore):
        self.store = page_store

    def match(self, outline: ContentOutline) -> AssemblyPlan:
        plan = AssemblyPlan(
            harmony_mode=outline.harmony_mode,
            preferred_family=outline.preferred_family,
        )
        for i, content_page in enumerate(outline.pages):
            candidates = self._query_candidates(content_page)
            if not candidates:
                candidates = self._degraded_query(content_page)
                if not candidates:
                    plan.warnings.append(f"P{i}: No pages found for type '{content_page.page_type.value}'")
                    continue

            scored = []
            for entry in candidates:
                score = self._score_entry(entry, content_page, outline.preferred_family)
                scored.append((score, entry))
            scored.sort(key=lambda x: x[0], reverse=True)
            best_score, best_entry = scored[0]
            fit = self._analyze_fit(best_entry, content_page)
            plan.pages.append(MatchResult(entry=best_entry, score=best_score, content_fit=fit))
        return plan

    def _query_candidates(self, content: PageContent) -> list[PageEntry]:
        family = content.preferred_family or None
        candidates = self.store.query(
            page_type=content.page_type, template_family=family,
            min_quality=0.3,
            has_image_slot=True if content.has_image else None,
            has_chart_slot=True if content.has_chart else None,
            limit=20,
        )
        if not candidates and family:
            candidates = self.store.query(page_type=content.page_type, min_quality=0.3, limit=20)
        return candidates

    def _degraded_query(self, content: PageContent) -> list[PageEntry]:
        similar = {
            PageType.CONTENT_2COL: [PageType.CONTENT_1COL, PageType.CONTENT_3COL],
            PageType.CONTENT_3COL: [PageType.CONTENT_2COL, PageType.CONTENT_1COL],
            PageType.IMAGE_LEFT: [PageType.IMAGE_RIGHT, PageType.IMAGE_FULL],
            PageType.IMAGE_RIGHT: [PageType.IMAGE_LEFT, PageType.IMAGE_FULL],
            PageType.SECTION_HEADER: [PageType.COVER, PageType.CONTENT_1COL],
            PageType.COVER: [PageType.SECTION_HEADER],
        }
        fallback_types = similar.get(content.page_type, [PageType.CONTENT_1COL])
        for ft in fallback_types:
            candidates = self.store.query(page_type=ft, min_quality=0.3, limit=10)
            if candidates:
                return candidates
        return self.store.query(min_quality=0.5, limit=10)

    def _score_entry(self, entry, content, preferred_family) -> float:
        score = 0.0
        if preferred_family and entry.template_family == preferred_family:
            score += 20
        if entry.page_type == content.page_type:
            score += 10
        else:
            score += 3

        title_len = len(content.title)
        if entry.content_constraints.max_title_chars > 0:
            ratio = title_len / max(entry.content_constraints.max_title_chars, 1)
            if ratio <= 1.0:
                score += 5
            elif ratio <= 1.5:
                score += 2
            else:
                score -= 5

        if content.has_image and entry.element_map.image_placeholders > 0:
            score += 5
        elif content.has_image and entry.element_map.image_placeholders == 0:
            score -= 2
        if content.has_chart and entry.element_map.chart_placeholders > 0:
            score += 5
        elif content.has_chart and entry.element_map.chart_placeholders == 0:
            score -= 2

        body_lines = content.body.count("\n") + 1 if content.body else 0
        if entry.content_constraints.max_body_lines > 0:
            cap_ratio = body_lines / max(entry.content_constraints.max_body_lines, 1)
            if cap_ratio <= 1.0:
                score += 3
            elif cap_ratio <= 1.5:
                score += 0
            else:
                score -= 3

        score += entry.quality_score * 5
        return score

    def _analyze_fit(self, entry, content) -> ContentFitResult:
        issues = []
        fit_map = {}
        score = 1.0

        title_slots = [s for s in entry.text_structure if s.role == "title"]
        if title_slots:
            fit_map["title"] = title_slots[0].shape_name
            if entry.content_constraints.max_title_chars > 0:
                if len(content.title) > entry.content_constraints.max_title_chars * 1.5:
                    score -= 0.3
                    issues.append(ContentIssue(field="title", message=f"标题过长 ({len(content.title)}/{entry.content_constraints.max_title_chars} chars)", severity=Severity.ERROR))
                elif len(content.title) > entry.content_constraints.max_title_chars:
                    score -= 0.1
                    issues.append(ContentIssue(field="title", message=f"标题略长 ({len(content.title)}/{entry.content_constraints.max_title_chars} chars)", severity=Severity.WARNING))

        subtitle_slots = [s for s in entry.text_structure if s.role == "subtitle"]
        if subtitle_slots and content.subtitle:
            fit_map["subtitle"] = subtitle_slots[0].shape_name

        body_slots = [s for s in entry.text_structure if s.role == "body"]
        if body_slots:
            fit_map["body"] = body_slots[0].shape_name
            body_lines = content.body.count("\n") + 1 if content.body else 0
            if entry.content_constraints.max_body_lines > 0:
                if body_lines > entry.content_constraints.max_body_lines * 1.5:
                    score -= 0.3
                    issues.append(ContentIssue(field="body", message=f"正文行数超出 ({body_lines}/{entry.content_constraints.max_body_lines} lines)", severity=Severity.ERROR))
                elif body_lines > entry.content_constraints.max_body_lines:
                    score -= 0.1
                    issues.append(ContentIssue(field="body", message=f"正文行数略多 ({body_lines}/{entry.content_constraints.max_body_lines} lines)", severity=Severity.WARNING))

        if content.has_image and entry.element_map.image_placeholders == 0:
            score -= 0.1
            issues.append(ContentIssue(field="image", message="需要图片但模板无图片位", severity=Severity.WARNING))
        if content.has_chart and entry.element_map.chart_placeholders == 0:
            score -= 0.1
            issues.append(ContentIssue(field="chart", message="需要图表但模板无图表位", severity=Severity.WARNING))

        return ContentFitResult(score=max(0.0, score), fit_map=fit_map, issues=issues)


class ContentAnalyzer:
    def analyze(self, entry, content) -> ContentFitResult:
        matcher = PageMatcher.__new__(PageMatcher)
        matcher.store = None
        return matcher._analyze_fit(entry, content)


class SlideAssembler:
    def __init__(self, harmony_mode=HarmonyMode.KEEP_SOURCE):
        self.harmony_mode = harmony_mode

    def assemble(self, plan, output_path) -> str:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        target_prs = Presentation()

        for i, match in enumerate(plan.pages):
            entry = match.entry
            source_path = entry.source_file
            if not source_path or not Path(source_path).exists():
                self._add_blank_slide(target_prs, f"[P{i}] Template not found: {entry.source_template}")
                continue

            try:
                source_prs = Presentation(str(source_path))
                if entry.slide_index >= len(source_prs.slides):
                    self._add_blank_slide(target_prs, f"[P{i}] Slide index out of range")
                    continue

                source_slide = source_prs.slides[entry.slide_index]
                target_slide_layout = target_prs.slide_layouts[0]
                target_slide = target_prs.slides.add_slide(target_slide_layout)
                self._clone_background(source_slide, target_slide)
                self._clone_shapes(source_slide, target_slide)
                self._fill_content(target_slide, match, plan)
            except Exception as e:
                self._add_blank_slide(target_prs, f"[P{i}] Error: {e}")

        if len(target_prs.slides) > 1:
            self._delete_slide(target_prs, 0)

        target_prs.save(str(output_path))
        return str(output_path)

    def _clone_background(self, source_slide, target_slide):
        try:
            layout = source_slide.slide_layout
            bg_xml = layout.background._element.xml if layout.background else ""
            match = re.search(r'srgbClr val="([^"]+)"', bg_xml)
            if match:
                color_hex = match.group(1)
                r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
                fill = target_slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r, g, b)
        except Exception:
            pass

    def _clone_shapes(self, source_slide, target_slide):
        sp_tree = target_slide.shapes._spTree
        for shape in source_slide.shapes:
            if shape.shape_type in (MSO_SHAPE_TYPE.SMART_ART, MSO_SHAPE_TYPE.GROUP):
                continue
            if shape.has_chart:
                continue
            try:
                new_element = deepcopy(shape._element)
                sp_tree.append(new_element)
            except Exception:
                continue

    def _fill_content(self, target_slide, match, plan):
        content = match.content_fit
        fit_map = content.fit_map
        for shape in target_slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_name = shape.name
            text_to_fill = None
            if "title" in fit_map and fit_map["title"] == shape_name:
                text_to_fill = getattr(match, '_content_title', "")
            elif "subtitle" in fit_map and fit_map["subtitle"] == shape_name:
                text_to_fill = getattr(match, '_content_subtitle', "")
            elif "body" in fit_map and fit_map["body"] == shape_name:
                text_to_fill = getattr(match, '_content_body', "")
            if text_to_fill:
                tf = shape.text_frame
                tf.clear()
                lines = text_to_fill.split("\n")
                for j, line in enumerate(lines):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line
                    try:
                        if p.runs:
                            p.runs[0].font.size = Pt(14)
                    except Exception:
                        pass

    def _add_blank_slide(self, prs, text):
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        try:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
            txBox.text_frame.text = text
        except Exception:
            pass

    def _delete_slide(self, prs, index):
        rId = prs.slides._sldIdLst[index].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[index]


class SlideValidator:
    def validate(self, output_path, plan) -> FillResult:
        output_path = Path(output_path)
        issues = []
        warnings = []
        stats = {}

        if not output_path.exists():
            return FillResult(output_path=str(output_path), passed=False, summary="Output file not found",
                             issues=[ValidationIssue(message="Output file not found", severity=Severity.ERROR)])

        try:
            prs = Presentation(str(output_path))
        except Exception as e:
            return FillResult(output_path=str(output_path), passed=False, summary=f"Cannot open output: {e}",
                             issues=[ValidationIssue(message=f"Cannot open: {e}", severity=Severity.ERROR)])

        actual_pages = len(prs.slides)
        expected_pages = len(plan.pages)

        if actual_pages != expected_pages:
            warnings.append(f"Page count: {actual_pages}/{expected_pages}")
            issues.append(ValidationIssue(check="page_count", message=f"Expected {expected_pages} pages, got {actual_pages}",
                                         severity=Severity.WARNING if actual_pages > 0 else Severity.ERROR))

        for i, slide in enumerate(prs.slides):
            shape_count = len(slide.shapes)
            if shape_count == 0:
                issues.append(ValidationIssue(page_index=i, check="shape_count", message="Empty slide (0 shapes)", severity=Severity.ERROR))
                continue
            if i < len(plan.pages):
                expected_shapes = (plan.pages[i].entry.element_map.text_boxes + plan.pages[i].entry.element_map.image_placeholders + plan.pages[i].entry.element_map.auto_shapes)
                if expected_shapes > 0 and shape_count < expected_shapes * 0.5:
                    issues.append(ValidationIssue(page_index=i, check="shape_count",
                                                 message=f"Shape count low: {shape_count}/{expected_shapes}", severity=Severity.WARNING))
            has_text = False
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    has_text = True
                    break
            if not has_text and shape_count > 0:
                issues.append(ValidationIssue(page_index=i, check="text_fill", message="No text content found on slide", severity=Severity.WARNING))

        for i, match in enumerate(plan.pages):
            for issue in match.content_fit.issues:
                if issue.severity == Severity.ERROR:
                    issues.append(ValidationIssue(page_index=i, check="content_overflow", message=issue.message, severity=Severity.ERROR))

        stats = {"total_pages": actual_pages, "expected_pages": expected_pages,
                 "errors": sum(1 for i in issues if i.severity == Severity.ERROR),
                 "warnings": sum(1 for i in issues if i.severity == Severity.WARNING)}

        passed = not any(i.severity == Severity.ERROR for i in issues)
        return FillResult(output_path=str(output_path), passed=passed,
                         summary=f"Validation: {stats['errors']} errors, {stats['warnings']} warnings",
                         issues=issues, warnings=warnings, stats=stats)


class TemplateFiller:
    def __init__(self, page_store):
        self.store = page_store
        self.matcher = PageMatcher(page_store)
        self.analyzer = ContentAnalyzer()
        self.assembler = SlideAssembler()
        self.validator = SlideValidator()

    def fill(self, outline, output_path, validate=True) -> FillResult:
        plan = self.matcher.match(outline)
        if not plan.pages:
            return FillResult(output_path=str(output_path), passed=False, summary="No matching pages found", warnings=plan.warnings)
        self.assembler.harmony_mode = outline.harmony_mode
        result_path = self.assembler.assemble(plan, output_path)
        if validate:
            result = self.validator.validate(result_path, plan)
            result.warnings.extend(plan.warnings)
            return result
        else:
            return FillResult(output_path=result_path, passed=True, summary="Assembly complete (validation skipped)", warnings=plan.warnings)
